import calendar
from sys import implementation
from django.shortcuts import render, redirect
from django.http import HttpResponse

import requests
import json

import pyrebase

import firebase_admin
from firebase_admin import auth
from firebase_admin import credentials
from firebase_admin import firestore

# import html to pdf converter
from django.template.loader import get_template
from xhtml2pdf import pisa
from django.contrib.staticfiles import finders
from io import BytesIO
from django.core.files import File
from django.conf import settings

import pytz
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches

import time

import os

from io import BytesIO
from pptx.dml.color import RGBColor

config = {
  'apiKey': "AIzaSyAh_oOakDenknpcWt5oucsLODSDiheWxps",
  'authDomain': "accreditation-management.firebaseapp.com",
  'projectId': "accreditation-management",
  'databaseURL': "https://accreditation-management-default-rtdb.firebaseio.com/",
  'storageBucket': "accreditation-management.appspot.com",
  'messagingSenderId': "553200895776",
  'appId': "1:553200895776:web:66ec853d3450869e4b3945"
}

firebase = pyrebase.initialize_app(config)
cred = credentials.Certificate("main_app/serviceAccountKey.json")
firebase_admin.initialize_app(cred)

auth_pyrebase = firebase.auth()

firestoreDB = firestore.client()

storage = firebase.storage()

# Create your views here.
def login(request):
    if 'user_id' in request.session:
        return redirect('/homepage')
    else:
        return render(request,'login.html')
    

def login_validation(request):
    if request.method == 'POST':
        try:
            email = request.POST.get('login_email')
            password = request.POST.get('login_password')

            user = auth_pyrebase.sign_in_with_email_and_password(email, password)

            request.session['user_id'] = user['localId']
            request.session['user_email'] = email

            it_department_accounts = firestoreDB.collection(u'it_department_accounts').document(str(user['localId']))
            hrm_department_accounts = firestoreDB.collection(u'hrm_department_accounts').document(str(user['localId']))
            it_department_doc = it_department_accounts.get()
            hrm_department_doc = hrm_department_accounts.get()
            if it_department_doc.exists:
                request.session['user_level'] = it_department_doc.to_dict()['user_level']
                request.session['middlename'] = it_department_doc.to_dict()['middlename']
                request.session['firstname'] = it_department_doc.to_dict()['firstname']
                request.session['lastname'] = it_department_doc.to_dict()['lastname']
            else:
                if hrm_department_doc.exists:
                    request.session['user_level'] = hrm_department_doc.to_dict()['user_level']
                    request.session['middlename'] = it_department_doc.to_dict()['middlename']
                    request.session['firstname'] = it_department_doc.to_dict()['firstname']
                    request.session['lastname'] = it_department_doc.to_dict()['lastname']
                else:
                    print('no document like this')
            return HttpResponse('Success!')
        except:
            return HttpResponse('Invalid Email or Password!')    

def homepage(request):
    if 'user_id' in request.session:
        return render(request, 'homepage.html')
    else:
        return redirect('/')

def storage_drive(request):
    return render(request, 'file_manager/storage_drive.html')

def upload_storage_drive(request):
    if request.method == 'POST':
        drive_upload =  request.FILES['drive_upload']

        selectLevel = request.POST.get('selectLevel')
        selectArea = request.POST.get('selectArea')
        selectParameter = request.POST.get('selectParameter')
        selectCategory = request.POST.get('selectCategory')
        selectDate = request.POST.get('selectDate')
        fileName = request.POST.get('fileName')

        uploadIn = request.POST.get('uploadIn')

        fileDirectory = selectLevel+"/"+selectArea+"/"+selectParameter+"/"+selectCategory+"/"+fileName

        doc_ref = firestoreDB.collection(selectLevel+'_'+selectArea+"_"+selectParameter+"_"+selectCategory).document()
        
        #upload to firebase storage
        storage.child(fileDirectory).put(drive_upload)

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(fileDirectory).get_url(None),
            'level': selectLevel,
            'area': selectArea,
            'parameter': selectParameter,
            'category': selectCategory,
            'date': selectDate,
            'file_name': fileName,
            'uploadIn': uploadIn,
        })

        tz = pytz.timezone('Asia/Hong_Kong')
        philippines_current_datetime = datetime.now(tz)

        doc_ref2 = firestoreDB.collection('activity_logs_storage_drive').document()

        doc_ref2.set({
            'activity_log_id': doc_ref2.id,
            'user_email': request.session['user_email'],
            'user_id': request.session['user_id'],
            'info': "uploaded " + fileName,
            'date': philippines_current_datetime,
            'uploadIn':uploadIn,
        })



        return redirect('storage_drive')

def activity_logs(request):
    if 'user_id' in request.session:
        logs = firestoreDB.collection('activity_logs_storage_drive').get()
        logs_data = []

        for log in logs:
            value = log.to_dict()
            logs_data.append(value)

        return render(request,'file_manager/activity_logs.html',{
        'logs_data': logs_data,
        })
    else:
        return redirect('/')

def recycle_bin(request):
    return render(request,'file_manager/recycle_bin.html')

#Level 1 / Area 1
def level1(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/level1.html')
    else:
        return redirect('/')
def area1(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area1/area1.html')
    else:
        return redirect('/')    

def level1_area1_parameterA(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 1_Parameter A_System').get()
        implementations = firestoreDB.collection('Level 1_Area 1_Parameter A_Implementation').get()
        outcomes = firestoreDB.collection('Level 1_Area 1_Parameter A_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area1_parameterA').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        
        return render(request,'file_manager/level1/area1/parameterA/parameterA.html', data)
    else:
        return redirect('/')

def level1_area1_parameterB(request):
    return render(request,'file_manager/level1/area1/parameterB/parameterB.html')

#Level 1 / Area 2
def area2(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area2/area2.html')
    else:
        return redirect('/')

def level1_area2_implementation(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area2/implementation.html')
    else:
        return redirect('/')

def level1_area2_implementation_parameterA(request):
    if 'user_id' in request.session:
        implementations = firestoreDB.collection('Level 1_Area 2_Parameter A_Implementation').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_implementation_parameterA').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/implementation/parameterA.html', data)
    else:
        return redirect('/')    

def level1_area2_implementation_parameterB(request):
    if 'user_id' in request.session:
        implementations = firestoreDB.collection('Level 1_Area 2_Parameter B_Implementation').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_implementation_parameterB').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/implementation/parameterB.html', data)
    else:
        return redirect('/')   

def level1_area2_implementation_parameterC(request):
    if 'user_id' in request.session:
        implementations = firestoreDB.collection('Level 1_Area 2_Parameter C_Implementation').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_implementation_parameterC').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/implementation/parameterC.html', data)
    else:
        return redirect('/')   

def level1_area2_implementation_parameterD(request):
    if 'user_id' in request.session:
        implementations = firestoreDB.collection('Level 1_Area 2_Parameter D_Implementation').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_implementation_parameterD').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/implementation/parameterD.html', data)
    else:
        return redirect('/')   

def level1_area2_implementation_parameterE(request):
    if 'user_id' in request.session:
        implementations = firestoreDB.collection('Level 1_Area 2_Parameter E_Implementation').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_implementation_parameterE').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/implementation/parameterE.html', data)
    else:
        return redirect('/')  

def level1_area2_implementation_parameterF(request):
    if 'user_id' in request.session:
        implementations = firestoreDB.collection('Level 1_Area 2_Parameter F_Implementation').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_implementation_parameterF').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/implementation/parameterF.html', data)
    else:
        return redirect('/') 

def level1_area2_implementation_parameterG(request):
    if 'user_id' in request.session:
        implementations = firestoreDB.collection('Level 1_Area 2_Parameter G_Implementation').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_implementation_parameterG').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/implementation/parameterG.html', data)
    else:
        return redirect('/') 

def level1_area2_implementation_parameterH(request):
    if 'user_id' in request.session:
        implementations = firestoreDB.collection('Level 1_Area 2_Parameter H_Implementation').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_implementation_parameterH').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/implementation/parameterH.html', data)
    else:
        return redirect('/') 

def level1_area2_outcome(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area2/outcome.html')
    else:
        return redirect('/')

def level1_area2_outcome_parameterA(request):
    if 'user_id' in request.session:
        outcomes = firestoreDB.collection('Level 1_Area 2_Parameter A_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_outcomes_parameterA').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/outcome/parameterA.html', data)
    else:
        return redirect('/')

def level1_area2_outcome_parameterB(request):
    if 'user_id' in request.session:
        outcomes = firestoreDB.collection('Level 1_Area 2_Parameter B_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_outcomes_parameterB').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/outcome/parameterB.html', data)
    else:
        return redirect('/')

def level1_area2_outcome_parameterC(request):
    if 'user_id' in request.session:
        outcomes = firestoreDB.collection('Level 1_Area 2_Parameter C_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_outcomes_parameterC').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/outcome/parameterC.html', data)
    else:
        return redirect('/')

def level1_area2_outcome_parameterD(request):
    if 'user_id' in request.session:
        outcomes = firestoreDB.collection('Level 1_Area 2_Parameter D_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_outcomes_parameterD').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/outcome/parameterD.html', data)
    else:
        return redirect('/')

def level1_area2_outcome_parameterE(request):
    if 'user_id' in request.session:
        outcomes = firestoreDB.collection('Level 1_Area 2_Parameter E_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_outcomes_parameterE').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/outcome/parameterE.html', data)
    else:
        return redirect('/')

def level1_area2_outcome_parameterF(request):
    if 'user_id' in request.session:
        outcomes = firestoreDB.collection('Level 1_Area 2_Parameter F_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_outcomes_parameterF').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/outcome/parameterF.html', data)
    else:
        return redirect('/')

def level1_area2_outcome_parameterG(request):
    if 'user_id' in request.session:
        outcomes = firestoreDB.collection('Level 1_Area 2_Parameter G_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_outcomes_parameterG').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/outcome/parameterG.html', data)
    else:
        return redirect('/')

def level1_area2_outcome_parameterH(request):
    if 'user_id' in request.session:
        outcomes = firestoreDB.collection('Level 1_Area 2_Parameter H_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_outcomes_parameterH').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/outcome/parameterH.html' , data)
    else:
        return redirect('/')

def level1_area2_system(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area2/system.html')
    else:
        return redirect('/')

def level1_area2_system_parameterA(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 2_Parameter A_System').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_system_parameterA').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/system/parameterA.html', data)
    else:
        return redirect('/')

def level1_area2_system_parameterB(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 2_Parameter B_System').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_system_parameterB').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/system/parameterB.html', data)
    else:
        return redirect('/')
        
def level1_area2_system_parameterC(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 2_Parameter C_System').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_system_parameterC').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/system/parameterC.html', data)
    else:
        return redirect('/')

def level1_area2_system_parameterD(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 2_Parameter D_System').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_system_parameterD').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/system/parameterD.html' , data)
    else:
        return redirect('/')

def level1_area2_system_parameterE(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 2_Parameter E_System').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_system_parameterE').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/system/parameterE.html', data)
    else:
        return redirect('/')

def level1_area2_system_parameterF(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 2_Parameter F_System').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_system_parameterF').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/system/parameterF.html', data)
    else:
        return redirect('/')

def level1_area2_system_parameterG(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 2_Parameter G_System').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_system_parameterG').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/system/parameterG.html', data)
    else:
        return redirect('/')

def level1_area2_system_parameterH(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 2_Parameter H_System').get()

        powerpoints = firestoreDB.collection('generatelevel1_area2_system_parameterH').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area2/system/parameterH.html', data)
    else:
        return redirect('/')

#Level 1 / Area 3
def area3(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area3/area3.html')
    else:
        return redirect('/')

def level1_area3_parameterA(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 3_Parameter A_System').get()
        implementations = firestoreDB.collection('Level 1_Area 3_Parameter A_Implementation').get()
        outcomes = firestoreDB.collection('Level 1_Area 3_Parameter A_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area3_parameterA').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area3/parameterA.html', data)
    else:
        return redirect('/')
        
def level1_area3_parameterB(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 3_Parameter B_System').get()
        implementations = firestoreDB.collection('Level 1_Area 3_Parameter B_Implementation').get()
        outcomes = firestoreDB.collection('Level 1_Area 3_Parameter B_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area3_parameterB').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area3/parameterB.html', data)
    else:
        return redirect('/')


def level1_area3_parameterC(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 3_Parameter C_System').get()
        implementations = firestoreDB.collection('Level 1_Area 3_Parameter C_Implementation').get()
        outcomes = firestoreDB.collection('Level 1_Area 3_Parameter C_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area3_parameterC').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area3/parameterC.html', data)
    else:
        return redirect('/')

def level1_area3_parameterD(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 3_Parameter D_System').get()
        implementations = firestoreDB.collection('Level 1_Area 3_Parameter D_Implementation').get()
        outcomes = firestoreDB.collection('Level 1_Area 3_Parameter D_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area3_parameterD').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area3/parameterD.html', data)
    else:
        return redirect('/')

def level1_area3_parameterE(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 3_Parameter E_System').get()
        implementations = firestoreDB.collection('Level 1_Area 3_Parameter E_Implementation').get()
        outcomes = firestoreDB.collection('Level 1_Area 3_Parameter E_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area3_parameterE').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area3/parameterE.html', data)
    else:
        return redirect('/')

def level1_area3_parameterF(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 3_Parameter F_System').get()
        implementations = firestoreDB.collection('Level 1_Area 3_Parameter F_Implementation').get()
        outcomes = firestoreDB.collection('Level 1_Area 3_Parameter F_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area3_parameterF').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])

        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area3/parameterF.html', data)
    else:
        return redirect('/')

#Level 1 / Area 4
def area4(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area4/area4.html')
    else:
        return redirect('/')
        
def level1_area4_implementation(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area4/implementation.html')
    else:
        return redirect('/')


def level1_area4_implementation_parameterA(request):
    if 'user_id' in request.session:
        implementations = firestoreDB.collection('Level 1_Area 4_Parameter A_Implementation').get()

        powerpoints = firestoreDB.collection('generatelevel1_area4_implementation_parameterA').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area4/implementation/parameterA.html', data)
    else:
        return redirect('/')

def level1_area4_implementation_parameterB(request):
    if 'user_id' in request.session:
        implementations = firestoreDB.collection('Level 1_Area 4_Parameter B_Implementation').get()

        powerpoints = firestoreDB.collection('generatelevel1_area4_implementation_parameterB').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area4/implementation/parameterB.html', data)
    else:
        return redirect('/')

def level1_area4_implementation_parameterC(request):
    if 'user_id' in request.session:
        implementations = firestoreDB.collection('Level 1_Area 4_Parameter C_Implementation').get()

        powerpoints = firestoreDB.collection('generatelevel1_area4_implementation_parameterC').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area4/implementation/parameterC.html', data)
    else:
        return redirect('/')

def level1_area4_implementation_parameterD(request):
    if 'user_id' in request.session:
        implementations = firestoreDB.collection('Level 1_Area 4_Parameter D_Implementation').get()

        powerpoints = firestoreDB.collection('generatelevel1_area4_implementation_parameterD').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area4/implementation/parameterD.html', data)
    else:
        return redirect('/')

def level1_area4_implementation_parameterE(request):
    if 'user_id' in request.session:
        implementations = firestoreDB.collection('Level 1_Area 4_Parameter E_Implementation').get()

        powerpoints = firestoreDB.collection('generatelevel1_area4_implementation_parameterE').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for implementation in implementations:
            value = implementation.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area4/implementation/parameterE.html', data)
    else:
        return redirect('/')

def level1_area4_outcome(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area4/outcome.html')
    else:
        return redirect('/')
        
def level1_area4_outcome_parameterA(request):
    if 'user_id' in request.session:
        outcomes = firestoreDB.collection('Level 1_Area 4_Parameter A_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area4_outcomes_parameterA').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area4/outcome/parameterA.html', data)
    else:
        return redirect('/')

def level1_area4_outcome_parameterB(request):
    if 'user_id' in request.session:
        outcomes = firestoreDB.collection('Level 1_Area 4_Parameter B_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area4_outcomes_parameterB').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area4/outcome/parameterB.html', data)
    else:
        return redirect('/')

def level1_area4_outcome_parameterC(request):
    if 'user_id' in request.session:
        outcomes = firestoreDB.collection('Level 1_Area 4_Parameter C_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area4_outcomes_parameterC').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area4/outcome/parameterC.html', data)
    else:
        return redirect('/')

def level1_area4_outcome_parameterD(request):
    if 'user_id' in request.session:
        outcomes = firestoreDB.collection('Level 1_Area 4_Parameter D_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area4_outcomes_parameterD').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area4/outcome/parameterD.html', data)
    else:
        return redirect('/')

def level1_area4_outcome_parameterE(request):
    if 'user_id' in request.session:
        outcomes = firestoreDB.collection('Level 1_Area 4_Parameter E_Outcomes').get()

        powerpoints = firestoreDB.collection('generatelevel1_area4_outcomes_parameterE').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for outcome in outcomes:
            value = outcome.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area4/outcome/parameterE.html', data)
    else:
        return redirect('/')

def level1_area4_system(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area4/system.html')
    else:
        return redirect('/')

def level1_area4_system_parameterA(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 4_Parameter A_System').get()

        powerpoints = firestoreDB.collection('generatelevel1_area4_system_parameterA').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area4/system/parameterA.html', data)
    else:
        return redirect('/')
def level1_area4_system_parameterB(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 4_Parameter B_System').get()

        powerpoints = firestoreDB.collection('generatelevel1_area4_system_parameterB').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area4/system/parameterB.html', data)
    else:
        return redirect('/')
def level1_area4_system_parameterC(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 4_Parameter C_System').get()

        powerpoints = firestoreDB.collection('generatelevel1_area4_system_parameterC').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area4/system/parameterC.html', data)
    else:
        return redirect('/')
def level1_area4_system_parameterD(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 4_Parameter D_System').get()

        powerpoints = firestoreDB.collection('generatelevel1_area4_system_parameterD').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area4/system/parameterD.html', data)
    else:
        return redirect('/')
def level1_area4_system_parameterE(request):
    if 'user_id' in request.session:
        systems = firestoreDB.collection('Level 1_Area 4_Parameter E_System').get()

        powerpoints = firestoreDB.collection('generatelevel1_area4_system_parameterE').get()

        uploaded_data = []
        needed_data = []
        generated_data = []

        for system in systems:
            value = system.to_dict()
            uploaded_data.append(value)
            needed_data.append(value['uploadIn'])


        for powerpoint in powerpoints:
            value = powerpoint.to_dict()
            generated_data.append(value)
            
        data = {
            'uploaded_data': uploaded_data,
            'needed_datas': needed_data,
            'generated_data': generated_data,
        }
        return render(request,'file_manager/level1/area4/system/parameterE.html', data)
    else:
        return redirect('/')

#Level 1 / Area 5
def area5(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area5/area5.html')
    else:
        return redirect('/')
def level1_area5_implementation(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area5/implementation.html')
    else:
        return redirect('/')
def level1_area5_implementation_parameterA(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area5/implementation/parameterA.html')
    else:
        return redirect('/')    
def level1_area5_implementation_parameterB(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area5/implementation/parameterB.html')
    else:
        return redirect('/')
def level1_area5_implementation_parameterC(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area5/implementation/parameterC.html')
    else:
        return redirect('/')
def level1_area5_implementation_parameterD(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area5/implementation/parameterD.html')
    else:
        return redirect('/')
def level1_area5_outcome(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area5/outcome.html')
    else:
        return redirect('/')
def level1_area5_outcome_parameterA(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area5/outcome/parameterA.html')
    else:
        return redirect('/')
def level1_area5_outcome_parameterB(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area5/outcome/parameterB.html')
    else:
        return redirect('/')

def level1_area5_outcome_parameterC(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area5/outcome/parameterC.html')
    else:
        return redirect('/')
def level1_area5_outcome_parameterD(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area5/outcome/parameterD.html')
    else:
        return redirect('/')
def level1_area5_system(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area5/system.html')
    else:
        return redirect('/')

def level1_area5_system_parameterA(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area5/system/parameterA.html')
    else:
        return redirect('/')    
def level1_area5_system_parameterB(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area5/system/parameterB.html')
    else:
        return redirect('/')    
def level1_area5_system_parameterC(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area5/system/parameterC.html')
    else:
        return redirect('/')    
def level1_area5_system_parameterD(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area5/system/parameterD.html')
    else:
        return redirect('/') 

#Level 1 / Area 6
def area6(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area6/area6.html')
    else:
        return redirect('/')
def level1_area6_parameterA(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area6/parameterA.html')
    else:
        return redirect('/')    
def level1_area6_parameterB(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area6/parameterB.html')
    else:
        return redirect('/')
def level1_area6_parameterC(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area6/parameterC.html')
    else:
        return redirect('/')
def level1_area6_parameterD(request):
    if 'user_id' in request.session:
        return render(request,'file_manager/level1/area6/parameterD.html')
    else:
        return redirect('/')





def logout(request):
    try:
        del request.session['user_id']
    except:
        return redirect('/')
    return redirect('/')

def manage_accounts(request):
    if 'user_id' not in request.session:
        return redirect('/')
    else:
        if request.method == 'GET':
            department = request.GET.get('department')

            if department == 'IT':
                getCollection = firestoreDB.collection('it_department_accounts').get()

                collection_data = []

                for user in getCollection:
                    value = user.to_dict()
                    collection_data.append(value)

                data = {
                    'user_data': collection_data,
                    'department': "IT",
                }
                
                return render(request, 'manage_accounts.html', data)
            
            if department == 'HRM':
                getCollection = firestoreDB.collection('hrm_department_accounts').get()

                collection_data = []

                for user in getCollection:
                    value = user.to_dict()
                    collection_data.append(value)

                data = {
                    'user_data': collection_data,
                    'department': "HRM",
                }
                return render(request, 'manage_accounts.html', data)
            
            if department == 'TOURISM':
                getCollection = firestoreDB.collection('tourism_department_accounts').get()

                collection_data = []

                for user in getCollection:
                    value = user.to_dict()
                    collection_data.append(value)

                data = {
                    'user_data': collection_data,
                    'department': "TOURISM",
                }
                return render(request, 'manage_accounts.html', data)
            
            if department == 'EDUC':
                getCollection = firestoreDB.collection('educ_department_accounts').get()

                collection_data = []

                for user in getCollection:
                    value = user.to_dict()
                    collection_data.append(value)

                data = {
                    'user_data': collection_data,
                    'department': "EDUC",
                }
                return render(request, 'manage_accounts.html', data)

        department = request.POST.get('department')

        getCollection = firestoreDB.collection('it_department_accounts').get()

        collection_data = []

        for user in getCollection:
            value = user.to_dict()
            collection_data.append(value)

        data = {
            'user_data': collection_data,
            'department': "IT",
        }
        return render(request, 'manage_accounts.html', data)


def addAccount(request):
    if request.method == 'POST':
        access_rights = request.POST.get('access_rights')
        firstname = request.POST.get('firstname')
        middlename = request.POST.get('middlename')
        lastname = request.POST.get('lastname')
        email = request.POST.get('email')
        password = request.POST.get('password')
        confirm_password = request.POST.get('confirm_password')
        contact = request.POST.get('contact')
        address = request.POST.get('address')
        user_level = request.POST.get('user_level')
        birthdate = request.POST.get('birthdate')

        if password != confirm_password:
            return HttpResponse("Password Do Not Match!")
        else:
            try:
                #register email and password to firebase auth
                user = auth_pyrebase.create_user_with_email_and_password(email, password)

                if access_rights == 'IT':
                    doc_ref = firestoreDB.collection('it_department_accounts').document(user['localId'])
                elif access_rights == 'HRM':
                    doc_ref = firestoreDB.collection('hrm_department_accounts').document(user['localId'])
                elif access_rights == 'TOURISM':
                    doc_ref = firestoreDB.collection('tourism_department_accounts').document(user['localId'])    
                elif access_rights == 'EDUC':
                    doc_ref = firestoreDB.collection('educ_department_accounts').document(user['localId']) 
                

                doc_ref.set({
                    'user_id': doc_ref.id,
                    'firstname': firstname,
                    'middlename': middlename,
                    'lastname': lastname,
                    'email': email,
                    'contact': contact,
                    'address': address,
                    'user_level': user_level,
                    'birthdate': birthdate,
                    'access_rights': access_rights,

                })
                return HttpResponse("Success!")
            except requests.HTTPError as e:
                error_json = e.args[1]
                error = json.loads(error_json)['error']['message']
                if error == "EMAIL_EXISTS":
                    return HttpResponse('Email Already Exists!')
            


def editAccount(request):
    if request.method == 'POST':
        firstname = request.POST.get('firstname')
        middlename = request.POST.get('middlename')
        lastname = request.POST.get('lastname')
        email = request.POST.get('email')
        password = request.POST.get('password')
        confirm_password = request.POST.get('confirm_password')
        contact = request.POST.get('contact')
        address = request.POST.get('address')
        user_level = request.POST.get('user_level')
        birthdate = request.POST.get('birthdate')

        user_id = request.POST.get('user_id')
        access_rights = request.POST.get('access_rights')

        if password != confirm_password:
            return HttpResponse("Password Do Not Match!")
        else:
            try:

                if access_rights == 'IT':
                    doc_ref = firestoreDB.collection('it_department_accounts').document(user_id)
                elif access_rights == 'HRM':
                    doc_ref = firestoreDB.collection('hrm_department_accounts').document(user_id)
                elif access_rights == 'TOURISM':
                    doc_ref = firestoreDB.collection('tourism_department_accounts').document(user_id)    
                elif access_rights == 'EDUC':
                    doc_ref = firestoreDB.collection('educ_department_accounts').document(user_id) 
                

                doc_ref.update({
                    'user_id': doc_ref.id,
                    'firstname': firstname,
                    'middlename': middlename,
                    'lastname': lastname,
                    'email': email,
                    'contact': contact,
                    'address': address,
                    'user_level': user_level,
                    'birthdate': birthdate,
                    'access_rights': access_rights,

                })
                return HttpResponse("Success!")
            except requests.HTTPError as e:
                error_json = e.args[1]
                error = json.loads(error_json)['error']['message']
                if error == "EMAIL_EXISTS":
                    return HttpResponse('Email Already Exists!')


#GENERATE AREA 1 PARAMETER A
def generatelevel1_area1_parameterA(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/test.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 1_Parameter A_System').get()
        implementations = firestoreDB.collection('Level 1_Area 1_Parameter A_Implementation').get()
        outcomes = firestoreDB.collection('Level 1_Area 1_Parameter A_Outcomes').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%201%20parameter%20A.%20Vision%20misionpng.png?alt=media&token=a1260394-b05f-4eec-94cb-9723077a17cf"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%201%20parameter%20A%203rd%20page.png?alt=media&token=ff4b99bd-2d62-4c85-b963-67843140c5dc"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3rd PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%201%20parameter%20A%204th%20page.png?alt=media&token=1672c90f-f81b-4fcc-b1b1-cb42ec18812d"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE
        
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/test.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterA_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area1/parameterA/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/test.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area1_parameterA').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })


        return redirect('/level1/area1/parameterA')
    else:
        return redirect('/')


#GENERATE AREA 2 PARAMETER A IMPLEMENTATION
def generatelevel1_area2_implementation_parameterA(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_implementation_parameterA.pptx')
        except:
            print("no file found")

        implementations = firestoreDB.collection('Level 1_Area 2_Parameter A_Implementation').get()
        
        dynamic_images = []

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"IMPLEMENTATION" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "IMPLEMENTATION"
        # END "IMPLEMENTATION" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_implementation_parameterA.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterA_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/implementation/parameterA/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_implementation_parameterA.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_implementation_parameterA').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/implementation/parameterA')
    else:
        return redirect('/')

#GENERATE AREA 2 PARAMETER B IMPLEMENTATION
def generatelevel1_area2_implementation_parameterB(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_implementation_parameterB.pptx')
        except:
            print("no file found")

        implementations = firestoreDB.collection('Level 1_Area 2_Parameter B_Implementation').get()
        
        dynamic_images = []

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"IMPLEMENTATION" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "IMPLEMENTATION"
        # END "IMPLEMENTATION" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_implementation_parameterB.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterB_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/implementation/parameterB/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_implementation_parameterB.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_implementation_parameterB').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/implementation/parameterB')
    else:
        return redirect('/')


#GENERATE AREA 2 PARAMETER C IMPLEMENTATION
def generatelevel1_area2_implementation_parameterC(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_implementation_parameterC.pptx')
        except:
            print("no file found")

        implementations = firestoreDB.collection('Level 1_Area 2_Parameter C_Implementation').get()
        
        dynamic_images = []

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"IMPLEMENTATION" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "IMPLEMENTATION"
        # END "IMPLEMENTATION" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_implementation_parameterC.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterC_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/implementation/parameterC/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_implementation_parameterC.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_implementation_parameterC').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/implementation/parameterC')
    else:
        return redirect('/')


#GENERATE AREA 2 PARAMETER D IMPLEMENTATION
def generatelevel1_area2_implementation_parameterD(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_implementation_parameterD.pptx')
        except:
            print("no file found")

        implementations = firestoreDB.collection('Level 1_Area 2_Parameter D_Implementation').get()
        
        dynamic_images = []

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"IMPLEMENTATION" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "IMPLEMENTATION"
        # END "IMPLEMENTATION" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_implementation_parameterD.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterD_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/implementation/parameterD/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_implementation_parameterD.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_implementation_parameterD').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/implementation/parameterD')
    else:
        return redirect('/')


#GENERATE AREA 2 PARAMETER E IMPLEMENTATION
def generatelevel1_area2_implementation_parameterE(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_implementation_parameterE.pptx')
        except:
            print("no file found")

        implementations = firestoreDB.collection('Level 1_Area 2_Parameter E_Implementation').get()
        
        dynamic_images = []

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"IMPLEMENTATION" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "IMPLEMENTATION"
        # END "IMPLEMENTATION" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_implementation_parameterE.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterE_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/implementation/parameterE/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_implementation_parameterE.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_implementation_parameterE').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/implementation/parameterE')
    else:
        return redirect('/')


#GENERATE AREA 2 PARAMETER F IMPLEMENTATION
def generatelevel1_area2_implementation_parameterF(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_implementation_parameterF.pptx')
        except:
            print("no file found")

        implementations = firestoreDB.collection('Level 1_Area 2_Parameter F_Implementation').get()
        
        dynamic_images = []

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"IMPLEMENTATION" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "IMPLEMENTATION"
        # END "IMPLEMENTATION" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_implementation_parameterF.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterF_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/implementation/parameterF/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_implementation_parameterF.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_implementation_parameterF').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/implementation/parameterF')
    else:
        return redirect('/')


#GENERATE AREA 2 PARAMETER G IMPLEMENTATION
def generatelevel1_area2_implementation_parameterG(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_implementation_parameterG.pptx')
        except:
            print("no file found")

        implementations = firestoreDB.collection('Level 1_Area 2_Parameter G_Implementation').get()
        
        dynamic_images = []

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"IMPLEMENTATION" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "IMPLEMENTATION"
        # END "IMPLEMENTATION" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_implementation_parameterG.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterG_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/implementation/parameterG/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_implementation_parameterG.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_implementation_parameterG').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/implementation/parameterG')
    else:
        return redirect('/')

#GENERATE AREA 2 PARAMETER H IMPLEMENTATION
def generatelevel1_area2_implementation_parameterH(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_implementation_parameterH.pptx')
        except:
            print("no file found")

        implementations = firestoreDB.collection('Level 1_Area 2_Parameter H_Implementation').get()
        
        dynamic_images = []

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"IMPLEMENTATION" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "IMPLEMENTATION"
        # END "IMPLEMENTATION" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_implementation_parameterH.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterH_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/implementation/parameterH/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_implementation_parameterH.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_implementation_parameterH').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/implementation/parameterH')
    else:
        return redirect('/')



#GENERATE AREA 2 PARAMETER A OUTCOMES
def generatelevel1_area2_outcomes_parameterA(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_outcomes_parameterA.pptx')
        except:
            print("no file found")

        outcomes = firestoreDB.collection('Level 1_Area 2_Parameter A_Outcomes').get()
        
        dynamic_images = []

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"OUTCOMES" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "OUTCOMES"
        # END "OUTCOMES" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_outcomes_parameterA.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterA_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/outcomes/parameterA/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_outcomes_parameterA.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_outcomes_parameterA').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/outcome/parameterA')
    else:
        return redirect('/')


#GENERATE AREA 2 PARAMETER B OUTCOMES
def generatelevel1_area2_outcomes_parameterB(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_outcomes_parameterB.pptx')
        except:
            print("no file found")

        outcomes = firestoreDB.collection('Level 1_Area 2_Parameter B_Outcomes').get()
        
        dynamic_images = []

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"OUTCOMES" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "OUTCOMES"
        # END "OUTCOMES" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_outcomes_parameterB.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterB_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/outcomes/parameterB/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_outcomes_parameterB.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_outcomes_parameterB').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/outcome/parameterB')
    else:
        return redirect('/')

#GENERATE AREA 2 PARAMETER C OUTCOMES
def generatelevel1_area2_outcomes_parameterC(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_outcomes_parameterC.pptx')
        except:
            print("no file found")

        outcomes = firestoreDB.collection('Level 1_Area 2_Parameter C_Outcomes').get()
        
        dynamic_images = []

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"OUTCOMES" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "OUTCOMES"
        # END "OUTCOMES" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_outcomes_parameterC.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterC_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/outcomes/parameterC/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_outcomes_parameterC.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_outcomes_parameterC').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/outcome/parameterC')
    else:
        return redirect('/')

#GENERATE AREA 2 PARAMETER D OUTCOMES
def generatelevel1_area2_outcomes_parameterD(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_outcomes_parameterD.pptx')
        except:
            print("no file found")

        outcomes = firestoreDB.collection('Level 1_Area 2_Parameter D_Outcomes').get()
        
        dynamic_images = []

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"OUTCOMES" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "OUTCOMES"
        # END "OUTCOMES" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_outcomes_parameterD.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterD_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/outcomes/parameterD/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_outcomes_parameterD.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_outcomes_parameterD').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/outcome/parameterD')
    else:
        return redirect('/')

#GENERATE AREA 2 PARAMETER E OUTCOMES
def generatelevel1_area2_outcomes_parameterE(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_outcomes_parameterE.pptx')
        except:
            print("no file found")

        outcomes = firestoreDB.collection('Level 1_Area 2_Parameter E_Outcomes').get()
        
        dynamic_images = []

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"OUTCOMES" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "OUTCOMES"
        # END "OUTCOMES" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_outcomes_parameterE.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterE_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/outcomes/parameterE/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_outcomes_parameterE.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_outcomes_parameterE').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/outcome/parameterE')
    else:
        return redirect('/')

#GENERATE AREA 2 PARAMETER F OUTCOMES
def generatelevel1_area2_outcomes_parameterF(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_outcomes_parameterF.pptx')
        except:
            print("no file found")

        outcomes = firestoreDB.collection('Level 1_Area 2_Parameter F_Outcomes').get()
        
        dynamic_images = []

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"OUTCOMES" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "OUTCOMES"
        # END "OUTCOMES" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_outcomes_parameterF.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterF_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/outcomes/parameterF/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_outcomes_parameterF.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_outcomes_parameterF').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/outcome/parameterF')
    else:
        return redirect('/')

#GENERATE AREA 2 PARAMETER G OUTCOMES
def generatelevel1_area2_outcomes_parameterG(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_outcomes_parameterG.pptx')
        except:
            print("no file found")

        outcomes = firestoreDB.collection('Level 1_Area 2_Parameter G_Outcomes').get()
        
        dynamic_images = []

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"OUTCOMES" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "OUTCOMES"
        # END "OUTCOMES" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_outcomes_parameterG.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterG_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/outcomes/parameterG/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_outcomes_parameterG.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_outcomes_parameterG').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/outcome/parameterG')
    else:
        return redirect('/')

#GENERATE AREA 2 PARAMETER H OUTCOMES
def generatelevel1_area2_outcomes_parameterH(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_outcomes_parameterH.pptx')
        except:
            print("no file found")

        outcomes = firestoreDB.collection('Level 1_Area 2_Parameter H_Outcomes').get()
        
        dynamic_images = []

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"OUTCOMES" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "OUTCOMES"
        # END "OUTCOMES" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_outcomes_parameterH.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterH_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/outcomes/parameterH/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_outcomes_parameterH.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_outcomes_parameterH').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/outcome/parameterH')
    else:
        return redirect('/')



#GENERATE AREA 2 PARAMETER A SYSTEM
def generatelevel1_area2_system_parameterA(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_system_parameterA.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 2_Parameter A_System').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%201%20parameter%20A.%20Vision%20misionpng.png?alt=media&token=a1260394-b05f-4eec-94cb-9723077a17cf"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20A%203rd%20page.png?alt=media&token=322f6afc-a19a-4ea3-a393-8fbfe61146cc"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3RD PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20A%204TH%20page.png?alt=media&token=2a595cec-4e17-4b89-be78-e90ec3b1fb31"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE

        #FOR 5TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20A%205th%20page.png?alt=media&token=b6765896-ed11-48f4-8fe6-b4cc01a6fc5a"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 5TH PAGE

        #"SYSTEM" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(3)
        title_shape.text = "SYSTEM INPUT AND PROCESSES"
        # END "SYSTEM" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_system_parameterA.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterA_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/system/parameterA/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_system_parameterA.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_system_parameterA').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/system/parameterA')
    else:
        return redirect('/')


#GENERATE AREA 2 PARAMETER B SYSTEM
def generatelevel1_area2_system_parameterB(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_system_parameterB.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 2_Parameter B_System').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%201%20parameter%20A.%20Vision%20misionpng.png?alt=media&token=a1260394-b05f-4eec-94cb-9723077a17cf"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20B%203RD%20page.png?alt=media&token=9b8973f3-5c97-4cc8-974e-8c8c2d3cbb2f"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3RD PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20B%204th%20page.png?alt=media&token=8bf986cd-a21d-4591-b325-19f564553db3"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE

        #FOR 5TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20B%205th%20page.png?alt=media&token=528a6cdb-dcd9-4234-bd2e-b6a1f8efd203"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 5TH PAGE

        #"SYSTEM" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(3)
        title_shape.text = "SYSTEM INPUT AND PROCESSES"
        # END "SYSTEM" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_system_parameterB.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterB_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/system/parameterB/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_system_parameterB.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_system_parameterB').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/system/parameterB')
    else:
        return redirect('/')


#GENERATE AREA 2 PARAMETER C SYSTEM
def generatelevel1_area2_system_parameterC(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_system_parameterC.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 2_Parameter C_System').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%201%20parameter%20A.%20Vision%20misionpng.png?alt=media&token=a1260394-b05f-4eec-94cb-9723077a17cf"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20C%203rd%20page.png?alt=media&token=38fab29f-058e-4fd1-bae8-73aa50fbc6c7"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3RD PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20C%204th%20page.png?alt=media&token=74a7d50a-4cf6-4709-8cb6-d312d857e019"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE

        #FOR 5TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20C%205th%20page.png?alt=media&token=d8de1f0f-7e76-418f-aff7-779549a9d7fa"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 5TH PAGE

        #FOR 6TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20C%206th%20page.png?alt=media&token=7432a420-e247-4c8e-8ca6-5d4f7baef873"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 6TH PAGE

        #"SYSTEM" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(3)
        title_shape.text = "SYSTEM INPUT AND PROCESSES"
        # END "SYSTEM" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_system_parameterC.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterC_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/system/parameterC/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_system_parameterC.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_system_parameterC').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/system/parameterC')
    else:
        return redirect('/')

#GENERATE AREA 2 PARAMETER D SYSTEM
def generatelevel1_area2_system_parameterD(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_system_parameterD.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 2_Parameter D_System').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%201%20parameter%20A.%20Vision%20misionpng.png?alt=media&token=a1260394-b05f-4eec-94cb-9723077a17cf"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20D%203rd%20page.png?alt=media&token=bfd2667b-2dbb-4b9d-8ad2-f0329001743c"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3RD PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20D%204TH%20page.png?alt=media&token=f0d46cb0-e261-4be6-a881-429f141eb692"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE

        #FOR 5TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20D%205TH%20page.png?alt=media&token=94580adb-955a-4c5d-a503-af5c0c8f1d45"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 5TH PAGE

        #"SYSTEM" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(3)
        title_shape.text = "SYSTEM INPUT AND PROCESSES"
        # END "SYSTEM" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_system_parameterD.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterD_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/system/parameterD/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_system_parameterD.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_system_parameterD').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/system/parameterD')
    else:
        return redirect('/')


#GENERATE AREA 2 PARAMETER E SYSTEM
def generatelevel1_area2_system_parameterE(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_system_parameterE.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 2_Parameter E_System').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%201%20parameter%20A.%20Vision%20misionpng.png?alt=media&token=a1260394-b05f-4eec-94cb-9723077a17cf"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20E%203rd%20page.png?alt=media&token=90e2d9ba-dc9e-469b-8519-51c6232fc6be"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3RD PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20E%204th%20page.png?alt=media&token=5b61769a-7939-45e0-b763-6baac00dea46"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE

        #FOR 5TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20E%205th%20page.png?alt=media&token=aeb9bf8a-06d4-45e0-ab76-a1dfae7a54dc"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 5TH PAGE

        #FOR 6TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20E%206th%20page.png?alt=media&token=f50672b6-822b-4ae1-bc85-b62d8b3d98fc"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 6TH PAGE

        #FOR 7TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20E%207th%20page.png?alt=media&token=645c5722-2010-4f0d-8e19-c306fe584c3c"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 7TH PAGE

        #"SYSTEM" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(3)
        title_shape.text = "SYSTEM INPUT AND PROCESSES"
        # END "SYSTEM" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_system_parameterE.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterE_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/system/parameterE/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_system_parameterE.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_system_parameterE').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/system/parameterE')
    else:
        return redirect('/')

#GENERATE AREA 2 PARAMETER F SYSTEM
def generatelevel1_area2_system_parameterF(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_system_parameterF.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 2_Parameter F_System').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%201%20parameter%20A.%20Vision%20misionpng.png?alt=media&token=a1260394-b05f-4eec-94cb-9723077a17cf"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20F%203rd%20page.png?alt=media&token=9a5048ed-242b-4aea-9f9d-619f237f2537"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3RD PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20F%204TH%20page.png?alt=media&token=c683d082-9d76-4b11-b183-30e4034027db"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE

        #FOR 5TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20F%205TH%20page.png?alt=media&token=1951013f-2e99-4554-8dd5-4910488f360e"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 5TH PAGE

        #FOR 6TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20F%206TH%20page.png?alt=media&token=0a6c9e03-c97f-4f66-a931-7e544099d104"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 6TH PAGE

        #"SYSTEM" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(3)
        title_shape.text = "SYSTEM INPUT AND PROCESSES"
        # END "SYSTEM" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_system_parameterF.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterF_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/system/parameterF/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_system_parameterF.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_system_parameterF').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/system/parameterF')
    else:
        return redirect('/')

#GENERATE AREA 2 PARAMETER G SYSTEM
def generatelevel1_area2_system_parameterG(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_system_parameterG.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 2_Parameter G_System').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%201%20parameter%20A.%20Vision%20misionpng.png?alt=media&token=a1260394-b05f-4eec-94cb-9723077a17cf"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20G%203rd%20page.png?alt=media&token=cd57f491-a5a0-4de6-9fa9-17fba9101f49"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3RD PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20G%204th%20page.png?alt=media&token=cb500ce1-7dd0-419f-8a12-61bbd5950160"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE

        #FOR 5TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20G%205th%20page.png?alt=media&token=fa58e4c9-bbbe-4281-8111-24d2a8a8cf4f"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 5TH PAGE

        #FOR 6TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20G%206th%20page.png?alt=media&token=c1b419b8-ed38-429f-bbd8-01210bc8c307"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 6TH PAGE

        #"SYSTEM" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(3)
        title_shape.text = "SYSTEM INPUT AND PROCESSES"
        # END "SYSTEM" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_system_parameterG.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterG_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/system/parameterG/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_system_parameterG.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_system_parameterG').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/system/parameterG')
    else:
        return redirect('/')

#GENERATE AREA 2 PARAMETER H SYSTEM
def generatelevel1_area2_system_parameterH(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area2_system_parameterH.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 2_Parameter H_System').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%201%20parameter%20A.%20Vision%20misionpng.png?alt=media&token=a1260394-b05f-4eec-94cb-9723077a17cf"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20H%203rd%20page.png?alt=media&token=d24eeb50-e764-4469-9db2-41d4f8f5993f"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3RD PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20H%204th%20page.png?alt=media&token=1ced432c-980b-4dd1-bd42-e29ba5207d95"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE

        #FOR 5TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20H%205th%20page.png?alt=media&token=501fa31d-7902-48c1-9991-2a96b604a250"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 5TH PAGE

        #FOR 6TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%202%20system%20parameter%20H%206th%20page.png?alt=media&token=e22ce164-3664-4156-9d8d-a0f971b6a277"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 6TH PAGE

        #"SYSTEM" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(3)
        title_shape.text = "SYSTEM INPUT AND PROCESSES"
        # END "SYSTEM" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area2_system_parameterH.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterH_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area2/system/parameterH/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area2_system_parameterH.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area2_system_parameterH').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area2/system/parameterH')
    else:
        return redirect('/')


#GENERATE AREA 3 PARAMETER A
def generatelevel1_area3_parameterA(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area3_parameterA.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 3_Parameter A_System').get()
        implementations = firestoreDB.collection('Level 1_Area 3_Parameter A_Implementation').get()
        outcomes = firestoreDB.collection('Level 1_Area 3_Parameter A_Outcomes').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2Farea%203%20parameter%20A%202nd%20page.png?alt=media&token=5fd778f8-5ee7-4329-9be1-b82f4d552842"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2Farea%203%20parameter%20A%203rd%20page.png?alt=media&token=7a7b21a1-3969-4935-88e7-6aa89877e79c"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3rd PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2Farea%203%20parameter%20A%204th%20page.png?alt=media&token=3a4e164b-f603-43b4-898e-6bc02fc3ee9c"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE
        
        #FOR 5TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2Farea%203%20parameter%20A%205th%20page.png?alt=media&token=529547f5-02c0-4353-b26f-152c850c5330"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 5TH PAGE

        #FOR 6TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2Farea%203%20parameter%20A%206th%20page.png?alt=media&token=bef5ae73-ff43-4ffb-aa41-4075f6b2ed69"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 6TH PAGE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area3_parameterA.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterA_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area3/parameterA/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area3_parameterA.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area3_parameterA').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })


        return redirect('/level1/area3/parameterA')
    else:
        return redirect('/')

#GENERATE AREA 3 PARAMETER B
def generatelevel1_area3_parameterB(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area3_parameterB.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 3_Parameter B_System').get()
        implementations = firestoreDB.collection('Level 1_Area 3_Parameter B_Implementation').get()
        outcomes = firestoreDB.collection('Level 1_Area 3_Parameter B_Outcomes').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2Farea%203%20parameter%20A%202nd%20page.png?alt=media&token=5fd778f8-5ee7-4329-9be1-b82f4d552842"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2FPARAMETER%20B%2Farea%203%20parameter%20B%203rd%20page.png?alt=media&token=ed9bdbc0-0193-42aa-8b9e-32f0d7e73128"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3rd PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2FPARAMETER%20B%2Farea%203%20parameter%20B%204TH%20page.png?alt=media&token=c3137da0-108f-47cd-9e6f-d7e97a133c3d"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE
        
        #FOR 5TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2FPARAMETER%20B%2Farea%203%20parameter%20B%205TH%20page.png?alt=media&token=e2eef323-d954-4d71-b94f-3e4e5ad3d020"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 5TH PAGE

        #FOR 6TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2FPARAMETER%20B%2Farea%203%20parameter%20B%206TH%20page.png?alt=media&token=d845e22c-d1f6-4313-b166-1f13219e6426"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 6TH PAGE

        #FOR 7TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2FPARAMETER%20B%2Farea%203%20parameter%20B%207TH%20page.png?alt=media&token=02f08b23-6843-47a3-b56a-8555b07df4f3"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 7TH PAGE

        #FOR 8TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2FPARAMETER%20B%2Farea%203%20parameter%20B%208TH%20page.png?alt=media&token=677688c7-b921-4e35-af4e-b087feca723c"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 8TH PAGE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area3_parameterB.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterB_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area3/parameterB/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area3_parameterB.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area3_parameterB').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })


        return redirect('/level1/area3/parameterB')
    else:
        return redirect('/')


#GENERATE AREA 3 PARAMETER C
def generatelevel1_area3_parameterC(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area3_parameterC.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 3_Parameter C_System').get()
        implementations = firestoreDB.collection('Level 1_Area 3_Parameter C_Implementation').get()
        outcomes = firestoreDB.collection('Level 1_Area 3_Parameter C_Outcomes').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2Farea%203%20parameter%20A%202nd%20page.png?alt=media&token=5fd778f8-5ee7-4329-9be1-b82f4d552842"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2FPARAMETER%20C%2Farea%203%20parameter%20C%203rd%20page.png?alt=media&token=0533da7b-3f9e-462f-89fd-7fccbe002fee"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3rd PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2FPARAMETER%20C%2Farea%203%20parameter%20C%204TH%20page.png?alt=media&token=74ae170d-e64a-438d-8ffc-f1ad39f8d622"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE
        
        #FOR 5TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2FPARAMETER%20C%2Farea%203%20parameter%20C%205TH%20page.png?alt=media&token=dc27c4ab-464b-473d-82ed-5f801825f92f"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 5TH PAGE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area3_parameterC.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterC_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area3/parameterC/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area3_parameterC.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area3_parameterC').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })


        return redirect('/level1/area3/parameterC')
    else:
        return redirect('/')

#GENERATE AREA 3 PARAMETER D
def generatelevel1_area3_parameterD(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area3_parameterD.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 3_Parameter D_System').get()
        implementations = firestoreDB.collection('Level 1_Area 3_Parameter D_Implementation').get()
        outcomes = firestoreDB.collection('Level 1_Area 3_Parameter D_Outcomes').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2Farea%203%20parameter%20A%202nd%20page.png?alt=media&token=5fd778f8-5ee7-4329-9be1-b82f4d552842"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2FPARAMETER%20D%2Farea%203%20parameter%20D%203rd%20page.png?alt=media&token=04a25653-0ce5-464f-b8c8-c7edddfc62f5"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3rd PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2FPARAMETER%20D%2Farea%203%20parameter%20D%204TH%20page.png?alt=media&token=857bd4a7-305c-402f-ab01-dc80d9619d22"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE
        
        #FOR 5TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2FPARAMETER%20D%2Farea%203%20parameter%20D%205TH%20page.png?alt=media&token=5465a55a-0bdd-4fb0-8bb1-5bf47e6b859c"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 5TH PAGE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area3_parameterD.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterD_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area3/parameterD/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area3_parameterD.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area3_parameterD').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })


        return redirect('/level1/area3/parameterD')
    else:
        return redirect('/')


#GENERATE AREA 3 PARAMETER E
def generatelevel1_area3_parameterE(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area3_parameterE.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 3_Parameter E_System').get()
        implementations = firestoreDB.collection('Level 1_Area 3_Parameter E_Implementation').get()
        outcomes = firestoreDB.collection('Level 1_Area 3_Parameter E_Outcomes').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%201%20parameter%20A.%20Vision%20misionpng.png?alt=media&token=a1260394-b05f-4eec-94cb-9723077a17cf"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2FPARAMETER%20E%2Farea%203%20parameter%20E%203rd%20page.png?alt=media&token=fcd10dd3-6536-4164-b127-83a46508a37f"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3rd PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2FPARAMETER%20E%2Farea%203%20parameter%20E%204TH%20page.png?alt=media&token=0c3c0970-2374-4edb-86b0-e49ef4592776"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE
           
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area3_parameterE.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterE_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area3/parameterE/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area3_parameterE.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area3_parameterE').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })


        return redirect('/level1/area3/parameterE')
    else:
        return redirect('/')

#GENERATE AREA 3 PARAMETER F
def generatelevel1_area3_parameterF(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area3_parameterF.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 3_Parameter F_System').get()
        implementations = firestoreDB.collection('Level 1_Area 3_Parameter F_Implementation').get()
        outcomes = firestoreDB.collection('Level 1_Area 3_Parameter F_Outcomes').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })

        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%201%20parameter%20A.%20Vision%20misionpng.png?alt=media&token=a1260394-b05f-4eec-94cb-9723077a17cf"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2FPARAMETER%20F%2Farea%203%20parameter%20F%203rd%20page.png?alt=media&token=490e1a35-12dd-4f83-8426-a258cdaeca2d"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3rd PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2FPARAMETER%20F%2Farea%203%20parameter%20F%204TH%20page.png?alt=media&token=96727aab-5a28-4a85-98e1-6a803fd2bb59"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE

        #FOR 5TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2FPARAMETER%20F%2Farea%203%20parameter%20F%205TH%20page.png?alt=media&token=422a581c-dca1-410c-979b-622b2e624f91"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 5TH PAGE

        #FOR 6TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2FPARAMETER%20F%2Farea%203%20parameter%20F%206TH%20page.png?alt=media&token=3a82ea59-a71e-48fa-9773-bff42c018b83"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 6TH PAGE

        #FOR 7TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%203%20PPT%20template%2FPARAMETER%20F%2Farea%203%20parameter%20F%207TH%20page.png?alt=media&token=c573ec73-cb50-4fef-a116-890fe5456503"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 7TH PAGE
           
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area3_parameterF.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterF_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area3/parameterF/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area3_parameterF.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area3_parameterF').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })


        return redirect('/level1/area3/parameterF')
    else:
        return redirect('/')


#GENERATE AREA 4 PARAMETER A IMPLEMENTATION
def generatelevel1_area4_implementation_parameterA(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area4_implementation_parameterA.pptx')
        except:
            print("no file found")

        implementations = firestoreDB.collection('Level 1_Area 4_Parameter A_Implementation').get()
        
        dynamic_images = []

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"IMPLEMENTATION" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "IMPLEMENTATION"
        # END "IMPLEMENTATION" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area4_implementation_parameterA.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterA_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area4/implementation/parameterA/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area4_implementation_parameterA.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area4_implementation_parameterA').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area4/implementation/parameterA')
    else:
        return redirect('/')

#GENERATE AREA 4 PARAMETER B IMPLEMENTATION
def generatelevel1_area4_implementation_parameterB(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area4_implementation_parameterB.pptx')
        except:
            print("no file found")

        implementations = firestoreDB.collection('Level 1_Area 4_Parameter B_Implementation').get()
        
        dynamic_images = []

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"IMPLEMENTATION" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "IMPLEMENTATION"
        # END "IMPLEMENTATION" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area4_implementation_parameterB.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterB_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area4/implementation/parameterB/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area4_implementation_parameterB.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area4_implementation_parameterB').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area4/implementation/parameterB')
    else:
        return redirect('/')

#GENERATE AREA 4 PARAMETER C IMPLEMENTATION
def generatelevel1_area4_implementation_parameterC(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area4_implementation_parameterC.pptx')
        except:
            print("no file found")

        implementations = firestoreDB.collection('Level 1_Area 4_Parameter C_Implementation').get()
        
        dynamic_images = []

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"IMPLEMENTATION" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "IMPLEMENTATION"
        # END "IMPLEMENTATION" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area4_implementation_parameterC.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterC_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area4/implementation/parameterC/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area4_implementation_parameterC.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area4_implementation_parameterC').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area4/implementation/parameterC')
    else:
        return redirect('/')

#GENERATE AREA 4 PARAMETER D IMPLEMENTATION
def generatelevel1_area4_implementation_parameterD(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area4_implementation_parameterD.pptx')
        except:
            print("no file found")

        implementations = firestoreDB.collection('Level 1_Area 4_Parameter D_Implementation').get()
        
        dynamic_images = []

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"IMPLEMENTATION" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "IMPLEMENTATION"
        # END "IMPLEMENTATION" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area4_implementation_parameterD.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterD_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area4/implementation/parameterD/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area4_implementation_parameterD.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area4_implementation_parameterD').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area4/implementation/parameterD')
    else:
        return redirect('/')

#GENERATE AREA 4 PARAMETER E IMPLEMENTATION
def generatelevel1_area4_implementation_parameterE(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area4_implementation_parameterE.pptx')
        except:
            print("no file found")

        implementations = firestoreDB.collection('Level 1_Area 4_Parameter E_Implementation').get()
        
        dynamic_images = []

        for implementation in implementations:
            value = implementation.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"IMPLEMENTATION" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "IMPLEMENTATION"
        # END "IMPLEMENTATION" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area4_implementation_parameterE.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterE_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area4/implementation/parameterE/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area4_implementation_parameterE.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area4_implementation_parameterE').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area4/implementation/parameterE')
    else:
        return redirect('/')


#GENERATE AREA 4 PARAMETER A OUTCOMES
def generatelevel1_area4_outcomes_parameterA(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area4_outcomes_parameterA.pptx')
        except:
            print("no file found")

        outcomes = firestoreDB.collection('Level 1_Area 4_Parameter A_Outcomes').get()
        
        dynamic_images = []

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"OUTCOMES" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "OUTCOMES"
        # END "OUTCOMES" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area4_outcomes_parameterA.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterA_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area4/outcomes/parameterA/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area4_outcomes_parameterA.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area4_outcomes_parameterA').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area4/outcome/parameterA')
    else:
        return redirect('/')

#GENERATE AREA 4 PARAMETER B OUTCOMES
def generatelevel1_area4_outcomes_parameterB(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area4_outcomes_parameterB.pptx')
        except:
            print("no file found")

        outcomes = firestoreDB.collection('Level 1_Area 4_Parameter B_Outcomes').get()
        
        dynamic_images = []

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"OUTCOMES" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "OUTCOMES"
        # END "OUTCOMES" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area4_outcomes_parameterB.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterB_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area4/outcomes/parameterB/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area4_outcomes_parameterB.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area4_outcomes_parameterB').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area4/outcome/parameterB')
    else:
        return redirect('/')

#GENERATE AREA 4 PARAMETER C OUTCOMES
def generatelevel1_area4_outcomes_parameterC(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area4_outcomes_parameterC.pptx')
        except:
            print("no file found")

        outcomes = firestoreDB.collection('Level 1_Area 4_Parameter C_Outcomes').get()
        
        dynamic_images = []

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"OUTCOMES" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "OUTCOMES"
        # END "OUTCOMES" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area4_outcomes_parameterC.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterC_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area4/outcomes/parameterC/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area4_outcomes_parameterC.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area4_outcomes_parameterC').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area4/outcome/parameterC')
    else:
        return redirect('/')

#GENERATE AREA 4 PARAMETER D OUTCOMES
def generatelevel1_area4_outcomes_parameterD(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area4_outcomes_parameterD.pptx')
        except:
            print("no file found")

        outcomes = firestoreDB.collection('Level 1_Area 4_Parameter D_Outcomes').get()
        
        dynamic_images = []

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"OUTCOMES" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "OUTCOMES"
        # END "OUTCOMES" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area4_outcomes_parameterD.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterD_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area4/outcomes/parameterD/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area4_outcomes_parameterD.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area4_outcomes_parameterD').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area4/outcome/parameterD')
    else:
        return redirect('/')

#GENERATE AREA 4 PARAMETER E OUTCOMES
def generatelevel1_area4_outcomes_parameterE(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area4_outcomes_parameterE.pptx')
        except:
            print("no file found")

        outcomes = firestoreDB.collection('Level 1_Area 4_Parameter E_Outcomes').get()
        
        dynamic_images = []

        for outcome in outcomes:
            value = outcome.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #"OUTCOMES" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(2.5)
        title_shape.text = "OUTCOMES"
        # END "OUTCOMES" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area4_outcomes_parameterE.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterE_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area4/outcomes/parameterE/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area4_outcomes_parameterE.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area4_outcomes_parameterE').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area4/outcome/parameterE')
    else:
        return redirect('/')


#GENERATE AREA 4 PARAMETER A SYSTEM
def generatelevel1_area4_system_parameterA(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area4_system_parameterA.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 4_Parameter A_System').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%201%20parameter%20A.%20Vision%20misionpng.png?alt=media&token=a1260394-b05f-4eec-94cb-9723077a17cf"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20A%2Farea%204%20system%20parameter%20A%203rd%20page.png?alt=media&token=8c7df844-304a-4e69-8739-7d11119f951a"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3RD PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20A%2Farea%204%20system%20parameter%20A%204TH%20page.png?alt=media&token=a0a8bd57-e11a-4d75-84bb-18d6b24c3252"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE

        #FOR 5TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20A%2Farea%204%20system%20parameter%20A%205TH%20page.png?alt=media&token=3e6d2f72-0a2a-4755-8f23-09e05ae8df00"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 5TH PAGE

        #FOR 6TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20A%2Farea%204%20system%20parameter%20A%206TH%20page.png?alt=media&token=34038bdf-3313-46ae-b5b1-5796b18ac3e8"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 6TH PAGE

        #FOR 7TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20A%2Farea%204%20system%20parameter%20A%207TH%20page.png?alt=media&token=788f92c4-bb80-4810-b97b-f464ad7bb951"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 7TH PAGE

        #"SYSTEM" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(3)
        title_shape.text = "SYSTEM INPUT AND PROCESSES"
        # END "SYSTEM" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area4_system_parameterA.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterA_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area4/system/parameterA/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area4_system_parameterA.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area4_system_parameterA').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area4/system/parameterA')
    else:
        return redirect('/')

#GENERATE AREA 4 PARAMETER B SYSTEM
def generatelevel1_area4_system_parameterB(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area4_system_parameterB.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 4_Parameter B_System').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%201%20parameter%20A.%20Vision%20misionpng.png?alt=media&token=a1260394-b05f-4eec-94cb-9723077a17cf"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20B%2Farea%204%20system%20parameter%20B%203rd%20page.png?alt=media&token=e4bd279b-3d36-46a0-a8d5-ee492fca3791"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3RD PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20B%2Farea%204%20system%20parameter%20B%204th%20page.png?alt=media&token=bfb750fe-9844-4948-ad3d-4a36ef3e0504"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE

        #FOR 5TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20B%2Farea%204%20system%20parameter%20B%205th%20page.png?alt=media&token=9c85a3b6-f354-478b-ad9a-d71c89f69d1a"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 5TH PAGE

        #FOR 6TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20B%2Farea%204%20system%20parameter%20B%206th%20page.png?alt=media&token=30d67a76-a340-400c-8591-a70a90138dd9"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 6TH PAGE

        #FOR 7TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20B%2Farea%204%20system%20parameter%20B%207TH%20page.png?alt=media&token=022c770f-fb31-40c7-82d0-82f8b900b14e"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 7TH PAGE

        #FOR 8TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20B%2Farea%204%20system%20parameter%20B%208TH%20page.png?alt=media&token=a1aee4f4-31f5-4ba5-ab8a-2f4ae4d3103d"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 8TH PAGE

        #FOR 9TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20B%2Farea%204%20system%20parameter%20B%209TH%20page.png?alt=media&token=c4d6a63b-c9a8-4fc7-bdbe-75b872a5b931"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 9TH PAGE

        #FOR 10TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20B%2Farea%204%20system%20parameter%20B%2010TH%20page.png?alt=media&token=0949ee17-24e4-4811-81ee-65f4e54fdd54"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 10TH PAGE

        #"SYSTEM" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(3)
        title_shape.text = "SYSTEM INPUT AND PROCESSES"
        # END "SYSTEM" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area4_system_parameterB.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterB_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area4/system/parameterB/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area4_system_parameterB.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area4_system_parameterB').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area4/system/parameterB')
    else:
        return redirect('/')


#GENERATE AREA 4 PARAMETER C SYSTEM
def generatelevel1_area4_system_parameterC(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area4_system_parameterC.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 4_Parameter C_System').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%201%20parameter%20A.%20Vision%20misionpng.png?alt=media&token=a1260394-b05f-4eec-94cb-9723077a17cf"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20C%2Farea%204%20system%20parameter%20C%203rd%20page.png?alt=media&token=0c05947c-121a-478b-9413-88bd042e45fb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3RD PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20C%2Farea%204%20system%20parameter%20C%204TH%20page.png?alt=media&token=064ff3d8-75c7-481a-8f5d-287ec87eabea"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE

        #FOR 5TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20C%2Farea%204%20system%20parameter%20C%205TH%20page.png?alt=media&token=35724abb-d7b9-411e-b61c-e5d3c31003d9"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 5TH PAGE

        #FOR 6TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20C%2Farea%204%20system%20parameter%20C%206TH%20page.png?alt=media&token=4dccd855-625c-41cf-bf66-17662b0ff3cb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 6TH PAGE

        #"SYSTEM" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(3)
        title_shape.text = "SYSTEM INPUT AND PROCESSES"
        # END "SYSTEM" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area4_system_parameterC.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterC_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area4/system/parameterC/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area4_system_parameterC.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area4_system_parameterC').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area4/system/parameterC')
    else:
        return redirect('/')

#GENERATE AREA 4 PARAMETER D SYSTEM
def generatelevel1_area4_system_parameterD(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area4_system_parameterD.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 4_Parameter D_System').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%201%20parameter%20A.%20Vision%20misionpng.png?alt=media&token=a1260394-b05f-4eec-94cb-9723077a17cf"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20D%2Farea%204%20system%20parameter%20D%203rd%20page.png?alt=media&token=f0694c73-3286-4382-aa58-a4dd0974ba41"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3RD PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20D%2Farea%204%20system%20parameter%20D%204th%20page.png?alt=media&token=73929a71-1057-4794-82aa-f3db371449df"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE

        #FOR 5TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20D%2Farea%204%20system%20parameter%20D%205th%20page.png?alt=media&token=d1cef1e5-caff-47ae-8d61-d86b4ca9193e"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 5TH PAGE

        #FOR 6TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20D%2Farea%204%20system%20parameter%20D%206th%20page.png?alt=media&token=5b50f59d-991b-4122-8d5a-7e6ef5bd00c7"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 6TH PAGE

        #FOR 7TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20D%2Farea%204%20system%20parameter%20D%207th%20page.png?alt=media&token=966f432b-cb03-4958-9b58-da2e50a01b17"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 7TH PAGE

        #FOR 8TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20D%2Farea%204%20system%20parameter%20D%208th%20page.png?alt=media&token=1f041888-bb4d-471c-b86b-e5bb8e56025d"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 8TH PAGE

        #FOR 9TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20D%2Farea%204%20system%20parameter%20D%209th%20page.png?alt=media&token=d9548cd5-173a-4e6a-b8d9-a7a4c9e5f217"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 9TH PAGE

        #"SYSTEM" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(3)
        title_shape.text = "SYSTEM INPUT AND PROCESSES"
        # END "SYSTEM" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area4_system_parameterD.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterD_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area4/system/parameterD/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area4_system_parameterD.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area4_system_parameterD').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area4/system/parameterD')
    else:
        return redirect('/')


#GENERATE AREA 4 PARAMETER E SYSTEM
def generatelevel1_area4_system_parameterE(request):
    if 'user_id' in request.session:
        try:
            os.remove('./ppt/level1_area4_system_parameterE.pptx')
        except:
            print("no file found")

        systems = firestoreDB.collection('Level 1_Area 4_Parameter E_System').get()
        
        dynamic_images = []

        for system in systems:
            value = system.to_dict()
            dynamic_images.append({
                'storage_file_url': value['storage_file_url'],
                'uploadIn': value['uploadIn'],
            })


        prs = Presentation()
        prs.slide_width = Inches(8.51)
        prs.slide_height = Inches(13.01)

        #0 = title Slide, 1 = title and content, 3 = section header, etc
        slide_layout = prs.slide_layouts[0]

        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/front_page.png?alt=media&token=4a37d5e1-e270-40fd-a51b-33d99689fefb"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)

        #FOR FRONT PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FRONT PAGE

        #FOR 2ND PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%201%20parameter%20A.%20Vision%20misionpng.png?alt=media&token=a1260394-b05f-4eec-94cb-9723077a17cf"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 2ND PAGE

        #FOR 3RD PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20E%2Farea%204%20system%20parameter%20E%203rd%20page.png?alt=media&token=c5891387-eedd-4cb4-b200-f6112a6d467d"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 3RD PAGE

        #FOR 4TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20E%2Farea%204%20system%20parameter%20E%204TH%20page.png?alt=media&token=75a0a40f-fa7c-43e6-8fa9-9d310535237e"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 4TH PAGE

        #FOR 5TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20E%2Farea%204%20system%20parameter%20E%205TH%20page.png?alt=media&token=ac1292a6-6078-4f1d-8791-c82f84187ab2"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 5TH PAGE

        #FOR 6TH PAGE
        #add front page slide
        slide = prs.slides.add_slide(slide_layout)
        #change background with an image of the slide …
        left = top = 0
        front_page_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/area%204%20PPT%20template%2FPARAMETER%20E%2Farea%204%20system%20parameter%20E%206TH%20page.png?alt=media&token=c6ff8f72-a98c-4da3-93b1-db5f1998510c"
        response_front_page = requests.get(front_page_img_url)
        image_data_front_page = BytesIO(response_front_page.content)
        front_page_pic = slide.shapes.add_picture(image_data_front_page, left-0.1*prs.slide_width, top, height = prs.slide_height)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 125, 95)
        #END FOR 6TH PAGE

        #"SYSTEM" TITLE SLIDE
        slide = prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(243, 241, 181)

        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.top = Inches(1)
        title_shape.left = Inches(0.1)
        title_shape.width = Inches(3)
        title_shape.text = "SYSTEM INPUT AND PROCESSES"
        # END "SYSTEM" TITLE SLIDE
        
        bulsu_img_url = "https://firebasestorage.googleapis.com/v0/b/accreditation-management.appspot.com/o/BULSU_logo.png?alt=media&token=10fc51f4-2689-468b-b7c4-e331d86540c1"
        response_bulsu = requests.get(bulsu_img_url)
        image_data_bulsu = BytesIO(response_bulsu.content)

        for images in dynamic_images:
            value = images
            slide = prs.slides.add_slide(slide_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(243, 241, 181)

            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.top = Inches(1)
            title_shape.left = Inches(0.1)
            title_shape.width = Inches(2)
            title_shape.text = value['uploadIn']

            #distance of the top edge
            top = Inches(2.5)
            #distance of the left edge
            left = Inches(0.3)
            height = Inches(5.5)
            
            img_url = value['storage_file_url']
            response = requests.get(img_url)
            image_data = BytesIO(response.content)

            pic = slide.shapes.add_picture(image_data, left, top, height=height)
            
            bulsuLogo = slide.shapes.add_picture(image_data_bulsu, Inches(0.3), Inches(11), height=Inches(1.5))

        prs.save('./ppt/level1_area4_system_parameterE.pptx')

        tz = pytz.timezone('Asia/Hong_Kong')
        now = datetime.now(tz)

        fileName = "parameterE_"+str(time.time())+".pptx"
        file_directory = "/ppt/level1/area4/system/parameterE/"+ fileName

        #upload image
        storage.child(file_directory).put('./ppt/level1_area4_system_parameterE.pptx')

        doc_ref = firestoreDB.collection('generatelevel1_area4_system_parameterE').document()

        doc_ref.set({
            'storage_file_id': doc_ref.id,
            'storage_file_url' : storage.child(file_directory).get_url(None),
            'file_name': fileName,
            'date': now,
        })
        return redirect('/level1/area4/system/parameterE')
    else:
        return redirect('/')










def feedbacks(request):
    feedbacks = firestoreDB.collection(u'feedbacks').get()
    return render(request,'feedback.html',{"feedbacks":[doc.to_dict() for doc in feedbacks]})
def todo_checklist(request):
    return render(request,'todo_checklist.html')
def send_feedback(request):
    content = request.GET.get('content')
    date_created = datetime.now()
    feedback_id =  calendar.timegm(date_created.timetuple())
    firestoreDB.collection(u'feedbacks').document(str(feedback_id)).set({
        'accreditor_name': request.session['firstname'] + " " + request.session['middlename'] + " " + request.session['lastname'],
        'date_created':date_created,
        'accreditor_id':request.session['user_id'],
        'feedback_id': str(feedback_id),
        'content': content,
    })
    return render(request,'homepage.html',{"validation":"Successfully Send your Feedback. Thanks for letting the admin know your side."})